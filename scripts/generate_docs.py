from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Pt
from docxcompose.composer import Composer
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Cm as PptCm
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parent.parent
DOWNLOADS = Path.home() / "Downloads"
OUTPUT_DIR = ROOT / "artifacts"
OUTPUT_DIR.mkdir(exist_ok=True)

TITLE_PATH = next(DOWNLOADS.glob("*Титульный*Лобин*docx"))
REPORT_TEMPLATE_PATH = next(DOWNLOADS.glob("*отчета и оформления.docx"))
PRESENTATION_TEMPLATE_PATH = next(DOWNLOADS.glob("*презентации*ДПО*pptx"))
HERO_PATH = ROOT / "frontend" / "src" / "assets" / "hero.png"


def remove_paragraph(paragraph) -> None:
    element = paragraph._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def remove_table(table) -> None:
    element = table._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def clear_document_content(document: Document) -> None:
    for table in list(document.tables):
        remove_table(table)
    for paragraph in list(document.paragraphs):
        remove_paragraph(paragraph)


def add_paragraph(document: Document, text: str, style: str = "No Spacing", bold: bool = False) -> None:
    paragraph = document.add_paragraph(style=style)
    run = paragraph.add_run(text)
    run.bold = bold
    run.font.name = "Times New Roman"
    run.font.size = Pt(14)
    if style in {"1ЗАГАЛОВОК", "2ЗАГАЛОВОК", "3ЗАГАЛОВОК"}:
        run.bold = True


def add_bullets(document: Document, items: list[str]) -> None:
    for item in items:
        add_paragraph(document, item, style="списокЧерта")


def add_sources(document: Document, items: list[str]) -> None:
    for item in items:
        add_paragraph(document, item, style="список")


def add_image_with_caption(document: Document, image_path: Path, caption: str) -> None:
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run()
    run.add_picture(str(image_path), width=Cm(14.5))

    caption_paragraph = document.add_paragraph(style="илюстрации")
    caption_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = caption_paragraph.add_run(caption)
    run.font.name = "Times New Roman"
    run.font.size = Pt(12)


def add_stack_table(document: Document) -> None:
    table = document.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    headers = ["Слой системы", "Технологии", "Назначение"]
    for idx, header in enumerate(headers):
        cell = table.rows[0].cells[idx]
        paragraph = cell.paragraphs[0]
        paragraph.style = "текстТаблицы"
        run = paragraph.add_run(header)
        run.bold = True

    rows = [
        ("Frontend", "React, TypeScript, Vite, Lucide React", "Интерфейс пользователя, трекинг привычек, отображение статистики"),
        ("Backend", "ASP.NET Core Web API, Entity Framework Core", "REST API, бизнес-логика, валидация и агрегация статистики"),
        ("Хранение данных", "PostgreSQL, Docker Compose", "Сохранение привычек, отметок выполнения и стартовых данных"),
    ]

    for row in rows:
        cells = table.add_row().cells
        for idx, value in enumerate(row):
            paragraph = cells[idx].paragraphs[0]
            paragraph.style = "текстТаблицы"
            paragraph.add_run(value)

    caption = document.add_paragraph(style="текстТаблицы")
    caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
    caption.add_run("Таблица 2.1 - Используемый технологический стек").bold = True


def build_report_body() -> Document:
    document = Document(str(REPORT_TEMPLATE_PATH))
    clear_document_content(document)

    for section in document.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(3)
        section.right_margin = Cm(1.5)

    add_paragraph(document, "ВВЕДЕНИЕ", style="1ЗАГАЛОВОК")
    add_paragraph(
        document,
        "В условиях высокой учебной и рабочей нагрузки пользователям сложно регулярно поддерживать полезные привычки, "
        "связанные с физической активностью, режимом сна, водным балансом и питанием. Даже при наличии мотивации человеку "
        "не всегда хватает наглядного инструмента, который помогает фиксировать прогресс и видеть динамику выполнения.",
    )
    add_paragraph(
        document,
        "В данной работе рассматривается процесс создания веб-приложения для трекинга привычек здорового образа жизни. "
        "Приложение позволяет добавлять привычки, отмечать их выполнение по дням, анализировать текущую серию, недельную "
        "активность и суммарную эффективность соблюдения намеченного режима.",
    )
    add_paragraph(
        document,
        "Актуальность предлагаемого решения обусловлена ростом интереса к цифровым инструментам самоконтроля и персональной "
        "аналитики. Пользователям нужны простые и понятные веб-сервисы, которые помогают переводить намерение в устойчивое действие.",
    )
    add_paragraph(
        document,
        "Целью работы является разработка веб-приложения для трекинга привычек здорового образа жизни на базе React, ASP.NET Core и PostgreSQL.",
    )
    add_paragraph(document, "Для достижения поставленной цели были решены следующие задачи:")
    add_bullets(
        document,
        [
            "проведен анализ конкурентных решений в области habit tracking;",
            "определены потребности целевой аудитории и ключевые пользовательские сценарии;",
            "спроектированы frontend-, backend- и database-компоненты системы;",
            "реализован прототип веб-приложения и подготовлены рекомендации по его развитию.",
        ],
    )

    add_paragraph(document, "1 АНАЛИТИЧЕСКИЙ РАЗДЕЛ", style="1ЗАГАЛОВОК")
    add_paragraph(
        document,
        "В данном разделе рассматриваются существующие цифровые решения для ведения привычек, особенности целевой аудитории, "
        "а также те проблемы, которые должна решать разрабатываемая система.",
    )

    add_paragraph(document, "1.1 Анализ рынка (Обзор конкурентов и существующих решений)", style="2ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Рынок приложений для самоконтроля и повышения личной эффективности насыщен мобильными и веб-решениями. "
        "Большинство из них строится вокруг ежедневного учета действий, визуализации серий и отправки напоминаний.",
    )

    add_paragraph(document, "1.1.1 Habitica", style="3ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Habitica использует игровую механику и превращает список привычек в RPG-интерфейс. Решение хорошо работает для "
        "пользователей, которым нужна высокая вовлеченность, однако может перегружать тех, кто предпочитает спокойный рабочий интерфейс.",
    )

    add_paragraph(document, "1.1.2 Loop Habit Tracker", style="3ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Loop Habit Tracker ориентирован на минималистичный учет привычек и отображение статистики. Основное достоинство решения "
        "заключается в простоте, однако веб-ориентированного сценария совместимого кроссплатформенного доступа в нем недостаточно.",
    )

    add_paragraph(document, "1.1.3 Streaks", style="3ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Streaks акцентирует внимание на сериях выполнений и визуальной чистоте интерфейса. Недостатком является зависимость "
        "от конкретной экосистемы устройств и сравнительно ограниченная гибкость серверной части для расширения аналитики.",
    )

    add_paragraph(document, "1.2 Анализ целевой аудитории", style="2ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Целевой аудиторией разрабатываемого веб-приложения являются студенты, молодые специалисты и пользователи, которые "
        "стремятся поддерживать базовые привычки здорового образа жизни без необходимости устанавливать сложные корпоративные сервисы.",
    )

    add_paragraph(document, "1.2.1 Описание целевой аудитории", style="3ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Для данной аудитории характерны переменный ритм дня, совмещение учебы и работы, частая работа за компьютером и потребность "
        "в простом цифровом помощнике. Пользователи ценят быстрый вход в систему, наглядную статистику и отсутствие перегруженного интерфейса.",
    )

    add_paragraph(document, "1.2.2 Проблемы целевой аудитории", style="3ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Основные проблемы аудитории выражаются в нерегулярности действий, слабой визуализации личного прогресса и отсутствии единого "
        "места, где можно быстро оценить динамику привычек за неделю. Дополнительным фактором выступает быстрая потеря мотивации, "
        "если пользователь не видит коротких, но понятных достижений.",
    )
    add_paragraph(
        document,
        "Предлагаемое решение закрывает эти проблемы за счет компактного интерфейса, отметок по дням, вычисления текущей серии, "
        "среднего процента выполнения и единого дашборда по последним семи дням.",
    )
    add_image_with_caption(document, HERO_PATH, "Рисунок 1.1 - Главный экран веб-приложения Healthy Habits")

    add_paragraph(document, "2 ТЕХНОЛОГИЧЕСКИЙ РАЗДЕЛ", style="1ЗАГАЛОВОК")
    add_paragraph(
        document,
        "В данном разделе определяются характеристики создаваемого приложения, а также требования к выбранному технологическому решению.",
    )
    add_paragraph(
        document,
        "Ключевыми требованиями выступают удобство использования, надежность хранения данных, возможность локального запуска в учебной среде "
        "и простота дальнейшего масштабирования функциональности.",
    )

    add_paragraph(document, "2.1 Характеристика веб-приложения", style="2ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Основным объектом предметной области является веб-приложение Healthy Habits, предназначенное для персонального учета привычек "
        "здорового образа жизни. Сфера применения решения включает индивидуальное использование в процессе самоорганизации и повседневного планирования.",
    )
    add_paragraph(
        document,
        "Назначение интерфейса заключается в том, чтобы предоставить пользователю единый экран для добавления привычек, отметки выполнения "
        "по датам, оценки недельного прогресса и просмотра агрегированной статистики без необходимости открывать несколько разделов.",
    )
    add_paragraph(document, "Основной функционал приложения включает:")
    add_bullets(
        document,
        [
            "создание привычек с указанием категории, цвета, иконки, недельной цели и времени напоминания;",
            "отметку выполнения привычки по конкретным датам с автоматическим пересчетом статистики;",
            "просмотр дашборда с показателями активности, выполнением за день и за неделю;",
            "удаление неактуальных привычек и начальную загрузку демонстрационных данных.",
        ],
    )

    add_paragraph(document, "2.2 Требования к технологическому решению", style="2ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Для реализации проекта выбрана клиент-серверная архитектура. Такой подход обеспечивает разделение ответственности между визуальным слоем, "
        "слоем прикладной логики и слоем хранения данных, а также упрощает последующее развитие продукта.",
    )
    add_paragraph(
        document,
        "В качестве frontend-решения использован React с TypeScript и сборщиком Vite. Выбор обусловлен высокой скоростью разработки, "
        "поддержкой компонентного подхода и удобной организацией состояния интерфейса.",
    )
    add_paragraph(
        document,
        "Backend реализован на ASP.NET Core Web API. Платформа позволяет быстро создавать REST API, подключать middleware, "
        "организовывать валидацию и интеграцию с ORM Entity Framework Core.",
    )
    add_paragraph(
        document,
        "Для хранения данных выбрана СУБД PostgreSQL. Она подходит для структурированного хранения привычек и отметок выполнения, "
        "поддерживает надежную работу с датами и легко запускается локально через Docker Compose.",
    )
    add_stack_table(document)

    add_paragraph(document, "3 ПРАКТИЧЕСКИЙ РАЗДЕЛ", style="1ЗАГАЛОВОК")
    add_paragraph(
        document,
        "В практической части описывается реализация интерфейса, серверной логики и тестирование созданного решения.",
    )

    add_paragraph(document, "3.1 Фронтенд-разработка", style="2ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Пользовательский интерфейс разработан на React с применением TypeScript. Основной экран состоит из блока метрик, полосы недельного прогресса, "
        "списка привычек и формы создания новой привычки. Для визуального акцента используется собственная цветовая схема, набор иконок Lucide и адаптивная сетка.",
    )
    add_paragraph(
        document,
        "Каждая привычка отображается отдельной карточкой, где пользователь видит категорию, целевой ритм, текущую серию и сетку последних семи дней. "
        "Нажатие на ячейку дня вызывает API и мгновенно обновляет состояние интерфейса.",
    )

    add_paragraph(document, "3.2 Бэкенд-разработка", style="2ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Серверная часть реализована в виде ASP.NET Core Minimal API. В приложении описаны маршруты для чтения списка привычек, "
        "создания новых записей, обновления данных, архивного удаления и переключения отметок выполнения.",
    )
    add_paragraph(
        document,
        "Для доступа к данным используется Entity Framework Core и контекст HabitsDbContext. На уровне модели реализованы сущности Habit и HabitCheckIn, "
        "а также ограничения целостности, включая уникальность отметки по связке привычка-дата.",
    )
    add_paragraph(
        document,
        "При запуске API выполняется инициализация базы данных и добавление демонстрационных привычек, что позволяет сразу открыть рабочий интерфейс "
        "без ручного наполнения таблиц.",
    )

    add_paragraph(document, "3.3 Тестирование, оценка и рекомендации", style="2ЗАГАЛОВОК")
    add_paragraph(
        document,
        "Проверка приложения выполнялась через сборку frontend-части и backend-проекта, а также через ручной сценарий использования: "
        "создание привычки, переключение отметки, удаление привычки и контроль пересчета статистики на главном экране.",
    )
    add_paragraph(
        document,
        "По результатам тестирования приложение корректно отображает активные привычки, рассчитывает число выполнений за день и неделю, "
        "формирует текущую серию и сохраняет данные в PostgreSQL.",
    )
    add_paragraph(document, "В качестве дальнейших рекомендаций можно выделить:")
    add_bullets(
        document,
        [
            "добавление авторизации и персональных аккаунтов пользователей;",
            "реализацию push- или email-напоминаний о привычках;",
            "расширение аналитики графиками за месяц и квартал;",
            "внедрение экспорта статистики и фильтрации по категориям.",
        ],
    )

    add_paragraph(document, "ЗАКЛЮЧЕНИЕ", style="1ЗАГАЛОВОК")
    add_paragraph(
        document,
        "В ходе выполнения работы было спроектировано и реализовано веб-приложение для трекинга привычек здорового образа жизни. "
        "Разработанное решение охватывает полный базовый цикл взаимодействия пользователя с системой: от создания привычки до анализа прогресса.",
    )
    add_paragraph(
        document,
        "Использование React, ASP.NET Core и PostgreSQL позволило создать современное, расширяемое и технологически целостное приложение. "
        "Полученный результат может применяться как самостоятельный учебный проект и как основа для дальнейшего развития полноценного сервиса habit tracking.",
    )

    add_paragraph(document, "СПИСОК ИСПОЛЬЗОВАННЫХ ИСТОЧНИКОВ", style="1ЗАГАЛОВОК")
    add_sources(
        document,
        [
            "Документация React [Электронный ресурс]. URL: https://react.dev/ (дата обращения: 20.06.2026).",
            "Документация ASP.NET Core [Электронный ресурс]. URL: https://learn.microsoft.com/aspnet/core/ (дата обращения: 20.06.2026).",
            "Документация PostgreSQL [Электронный ресурс]. URL: https://www.postgresql.org/docs/ (дата обращения: 20.06.2026).",
            "Документация Entity Framework Core [Электронный ресурс]. URL: https://learn.microsoft.com/ef/core/ (дата обращения: 20.06.2026).",
            "Документация Docker Compose [Электронный ресурс]. URL: https://docs.docker.com/compose/ (дата обращения: 20.06.2026).",
        ],
    )

    return document


def replace_shape_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    frame = shape.text_frame
    frame.clear()
    for index, part in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = part
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = PptPt(18)


def delete_pictures(slide) -> None:
    picture_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    for shape in picture_shapes:
        shape._element.getparent().remove(shape._element)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()
    for idx, part in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = part
        run.font.name = "Arial"
        run.font.size = PptPt(font_size)
        run.font.bold = bold
    return textbox


def build_presentation() -> Presentation:
    prs = Presentation(str(PRESENTATION_TEMPLATE_PATH))
    for slide in prs.slides:
        delete_pictures(slide)

    # Slide 1
    replace_shape_text(prs.slides[0].shapes[1], "Веб-приложение для трекинга привычек здорового образа жизни")
    replace_shape_text(prs.slides[0].shapes[2], "Выполнил: Лобин Александр Ильич\nГруппа: ЭФБО-06-24")
    prs.slides[0].shapes.add_picture(str(HERO_PATH), PptCm(18.3), PptCm(8.8), width=PptCm(10.2))

    # Slide 2
    slide = prs.slides[1]
    replace_shape_text(slide.shapes[0], "Цель и задачи работы")
    replace_shape_text(slide.shapes[4], "Цель")
    replace_shape_text(slide.shapes[5], "Разработать веб-приложение для удобного контроля привычек здорового образа жизни")
    replace_shape_text(slide.shapes[9], "Задачи")
    replace_shape_text(slide.shapes[10], "Проанализировать рынок habit tracking\nОпределить сценарии пользователей\nРеализовать frontend, backend и базу данных")
    replace_shape_text(slide.shapes[14], "Результат")
    replace_shape_text(slide.shapes[15], "Рабочее приложение с дашбордом, недельной сеткой отметок и статистикой")

    # Slide 3
    slide = prs.slides[2]
    replace_shape_text(slide.shapes[1], "Анализ конкурентов")
    cards = [
        ("01", "Habitica", "Сильная игровая мотивация, но перегруженный сценарий для спокойного трекинга"),
        ("02", "Loop Habit Tracker", "Минималистичный учет привычек, однако слабее выражен веб-сценарий доступа"),
        ("03", "Streaks", "Хорошая визуализация серий, но привязка к конкретной экосистеме устройств"),
        ("04", "Вывод", "Нужен легкий веб-интерфейс с быстрой отметкой привычек и понятной аналитикой"),
    ]
    mapping = [(2, 4, 5), (6, 8, 9), (10, 12, 13), (14, 16, 17)]
    for (num_idx, title_idx, desc_idx), (num, title, desc) in zip(mapping, cards):
        replace_shape_text(slide.shapes[num_idx], num)
        replace_shape_text(slide.shapes[title_idx], title)
        replace_shape_text(slide.shapes[desc_idx], desc)

    # Slide 4
    slide = prs.slides[3]
    replace_shape_text(slide.shapes[1], "Требования к веб-приложению")
    items = [
        ("Простота интерфейса", "Все основные действия доступны с одного экрана"),
        ("Надежность данных", "Отметки привычек сохраняются в PostgreSQL"),
        ("Быстрый запуск", "Локальный старт через Docker Compose, ASP.NET Core и Vite"),
        ("Расширяемость", "Архитектура позволяет добавить авторизацию, уведомления и расширенную аналитику"),
    ]
    mapping = [(4, 5), (8, 9), (12, 13), (16, 17)]
    for (title_idx, desc_idx), (title, desc) in zip(mapping, items):
        replace_shape_text(slide.shapes[title_idx], title)
        replace_shape_text(slide.shapes[desc_idx], desc)

    # Slide 5
    slide = prs.slides[4]
    replace_shape_text(slide.shapes[0], "Используемый стек")
    add_textbox(slide, PptCm(2.4), PptCm(5.0), PptCm(8.8), PptCm(8.2), "Frontend\nReact\nTypeScript\nVite\nLucide React", font_size=20, bold=True)
    add_textbox(slide, PptCm(12.0), PptCm(5.0), PptCm(8.8), PptCm(8.2), "Backend\nASP.NET Core Web API\nEntity Framework Core\nMinimal API", font_size=20, bold=True)
    add_textbox(slide, PptCm(21.5), PptCm(5.0), PptCm(8.8), PptCm(8.2), "Data\nPostgreSQL\nDocker Compose\nSeed data", font_size=20, bold=True)

    # Slide 6
    slide = prs.slides[5]
    replace_shape_text(slide.shapes[1], "Результаты разработки")
    slide.shapes.add_picture(str(HERO_PATH), PptCm(2.1), PptCm(4.4), width=PptCm(15.5))
    add_textbox(
        slide,
        PptCm(19.0),
        PptCm(5.0),
        PptCm(10.5),
        PptCm(8.5),
        "Реализовано:\n- дашборд с ключевыми метриками\n- карточки привычек с сеткой 7 дней\n- форма создания и удаления привычек\n- API и статистика по сериям\n\nЛокальный запуск:\nfrontend: http://localhost:5173\nbackend: http://localhost:5000",
        font_size=18,
    )

    # Slide 7
    slide = prs.slides[6]
    replace_shape_text(slide.shapes[1], "Архитектура решения")
    replace_shape_text(slide.shapes[3], "Клиент")
    replace_shape_text(slide.shapes[4], "React-интерфейс отправляет запросы к API и отображает привычки, сетку отметок и агрегированную статистику")
    replace_shape_text(slide.shapes[6], "Сервер")
    replace_shape_text(slide.shapes[7], "ASP.NET Core обрабатывает CRUD-операции, валидацию, расчет streak и weekly progress")
    replace_shape_text(slide.shapes[9], "База данных")
    replace_shape_text(slide.shapes[10], "PostgreSQL хранит сущности Habit и HabitCheckIn, а Docker Compose упрощает локальное развертывание")

    # Slide 8
    slide = prs.slides[7]
    replace_shape_text(slide.shapes[0], "Целевая аудитория")
    replace_shape_text(slide.shapes[2], "Студенты и молодые специалисты")
    replace_shape_text(slide.shapes[3], "Нужен простой способ не терять базовые привычки на фоне высокой нагрузки")
    replace_shape_text(slide.shapes[5], "Пользователи ЗОЖ-сервисов")
    replace_shape_text(slide.shapes[6], "Важно быстро отмечать воду, сон, активность и видеть результат без перегрузки интерфейса")
    replace_shape_text(slide.shapes[8], "Ключевая проблема")
    replace_shape_text(slide.shapes[9], "Отсутствие наглядной и регулярной обратной связи мешает превратить намерение в устойчивый режим")

    # Slide 9
    slide = prs.slides[8]
    replace_shape_text(slide.shapes[1], "Ключевые показатели решения")
    replace_shape_text(slide.shapes[2], "7")
    replace_shape_text(slide.shapes[3], "Дней в сетке привычки")
    replace_shape_text(slide.shapes[4], "100%")
    replace_shape_text(slide.shapes[5], "Локальный контроль данных")
    replace_shape_text(slide.shapes[7], "Dashboard")
    replace_shape_text(slide.shapes[8], "Активные привычки, выполнение за день и неделю, лучшая серия")
    replace_shape_text(slide.shapes[10], "API")
    replace_shape_text(slide.shapes[11], "CRUD привычек, toggle check-in, вычисление completion rate и streak")
    replace_shape_text(slide.shapes[13], "Docker + PostgreSQL")
    replace_shape_text(slide.shapes[14], "Быстрый локальный старт без конфликта с системной БД через порт 5433")

    # Slide 10
    slide = prs.slides[9]
    replace_shape_text(slide.shapes[0], "Итоги и развитие проекта")
    replace_shape_text(slide.shapes[2], "Плюсы")
    replace_shape_text(slide.shapes[3], "Понятный пользовательский сценарий\nТехнологически целостный стек\nНаглядная аналитика по привычкам")
    replace_shape_text(slide.shapes[5], "Дальнейшее развитие")
    replace_shape_text(slide.shapes[6], "Авторизация\nНапоминания\nФильтры и графики\nИстория по месяцам")
    replace_shape_text(slide.shapes[7], "Разработанное решение может использоваться как учебный проект и как база для дальнейшего развития полноценного сервиса habit tracking.")

    return prs


def compose_report() -> Path:
    body = build_report_body()
    body_path = OUTPUT_DIR / "report_body.docx"
    body.save(body_path)

    title_doc = Document(str(TITLE_PATH))
    body_doc = Document(str(body_path))
    composer = Composer(title_doc)
    composer.append(body_doc)

    result_path = OUTPUT_DIR / "Отчет_Healthy_Habits_Лобин_А.И..docx"
    composer.save(str(result_path))
    body_path.unlink(missing_ok=True)
    return result_path


def compose_presentation() -> Path:
    prs = build_presentation()
    result_path = OUTPUT_DIR / "Презентация_Healthy_Habits_Лобин_А.И..pptx"
    prs.save(str(result_path))
    return result_path


def main() -> None:
    report_path = compose_report()
    presentation_path = compose_presentation()
    print(report_path)
    print(presentation_path)


if __name__ == "__main__":
    main()
